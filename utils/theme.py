"""Theme management for PPT generation"""

import json
from pathlib import Path
from typing import Dict, Any
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class ThemeManager:
    """Manages PPT themes and styling"""

    DEFAULT_THEMES = {
        'default': {
            'name': 'Default',
            'background': (255, 255, 255),
            'title_color': (31, 78, 121),
            'text_color': (0, 0, 0),
            'accent_color': (68, 114, 196),
            'code_bg': (240, 240, 240),
            'code_text': (51, 51, 51),
            'font_title': 'Arial',
            'font_body': 'Arial',
            'font_code': 'Courier New',
            'title_size': 32,
            'heading_size': 24,
            'body_size': 14,
            'code_size': 11,
        },
        'dark': {
            'name': 'Dark',
            'background': (30, 30, 30),
            'title_color': (100, 200, 255),
            'text_color': (230, 230, 230),
            'accent_color': (0, 149, 255),
            'code_bg': (45, 45, 45),
            'code_text': (200, 200, 200),
            'font_title': 'Arial',
            'font_body': 'Arial',
            'font_code': 'Courier New',
            'title_size': 32,
            'heading_size': 24,
            'body_size': 14,
            'code_size': 11,
        },
        'light': {
            'name': 'Light',
            'background': (250, 250, 250),
            'title_color': (0, 102, 204),
            'text_color': (51, 51, 51),
            'accent_color': (255, 153, 0),
            'code_bg': (245, 245, 245),
            'code_text': (68, 68, 68),
            'font_title': 'Calibri',
            'font_body': 'Calibri',
            'font_code': 'Consolas',
            'title_size': 32,
            'heading_size': 24,
            'body_size': 14,
            'code_size': 11,
        },
        'professional': {
            'name': 'Professional',
            'background': (255, 255, 255),
            'title_color': (0, 32, 96),
            'text_color': (64, 64, 64),
            'accent_color': (192, 0, 0),
            'code_bg': (242, 242, 242),
            'code_text': (51, 51, 51),
            'font_title': 'Georgia',
            'font_body': 'Georgia',
            'font_code': 'Monaco',
            'title_size': 36,
            'heading_size': 26,
            'body_size': 16,
            'code_size': 12,
        }
    }

    def __init__(self, theme_name: str = 'default'):
        """Initialize theme manager

        Args:
            theme_name: Name of the theme to use
        """
        self.theme = self.load_theme(theme_name)

    def load_theme(self, theme_name: str) -> Dict[str, Any]:
        """Load theme configuration

        Args:
            theme_name: Name of the theme or path to custom theme file

        Returns:
            Theme configuration dictionary
        """
        # Check if it's a built-in theme
        if theme_name in self.DEFAULT_THEMES:
            return self.DEFAULT_THEMES[theme_name].copy()

        # Try to load from file
        theme_path = Path(theme_name)
        if theme_path.exists() and theme_path.suffix == '.json':
            with open(theme_path, 'r', encoding='utf-8') as f:
                return json.load(f)

        # Fallback to default
        print(f"Warning: Theme '{theme_name}' not found, using default theme")
        return self.DEFAULT_THEMES['default'].copy()

    def get_color(self, color_name: str) -> RGBColor:
        """Get RGB color from theme

        Args:
            color_name: Name of the color in theme

        Returns:
            RGBColor object
        """
        color = self.theme.get(color_name, (0, 0, 0))
        return RGBColor(*color)

    def get_font_size(self, size_name: str) -> int:
        """Get font size from theme

        Args:
            size_name: Name of the size in theme

        Returns:
            Font size in points
        """
        return self.theme.get(size_name, 14)

    def get_font_name(self, font_type: str) -> str:
        """Get font name from theme

        Args:
            font_type: Type of font (title, body, code)

        Returns:
            Font name
        """
        return self.theme.get(f'font_{font_type}', 'Arial')

    def apply_text_style(self, text_frame, style_type: str = 'body'):
        """Apply text styling to a text frame

        Args:
            text_frame: PPT text frame object
            style_type: Type of style (title, heading, body, code)
        """
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(self.get_font_size(f'{style_type}_size'))

            if style_type == 'title':
                paragraph.font.name = self.get_font_name('title')
                paragraph.font.color.rgb = self.get_color('title_color')
                paragraph.font.bold = True
            elif style_type == 'heading':
                paragraph.font.name = self.get_font_name('body')
                paragraph.font.color.rgb = self.get_color('title_color')
                paragraph.font.bold = True
            elif style_type == 'code':
                paragraph.font.name = self.get_font_name('code')
                paragraph.font.color.rgb = self.get_color('code_text')
            else:
                paragraph.font.name = self.get_font_name('body')
                paragraph.font.color.rgb = self.get_color('text_color')

    def save_theme(self, filepath: str):
        """Save current theme to file

        Args:
            filepath: Path to save theme configuration
        """
        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(self.theme, f, indent=2)
